from fastapi import FastAPI, Request
from fastapi import UploadFile, File, Form
from typing import List
from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Literal, Optional
from openai import OpenAI
from dotenv import load_dotenv
import os
import re
import json
import base64
import asyncio
import subprocess
import shutil
from tempfile import TemporaryDirectory
from pathlib import Path

try:
    from docx import Document
except ImportError:
    Document = None

load_dotenv()

client = OpenAI(
    api_key=os.getenv("OPENAI_API_KEY"),
    base_url=os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1"),
)

# 模型选择：在 .env 中设置 OPENAI_MODEL 可覆盖默认值。推荐见 模型推荐.md
MODEL_NAME = os.getenv("OPENAI_MODEL", "gpt-4o")

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


class MeetingInfo(BaseModel):
    time: str
    topic: str
    roles: List[str]


class GenerateRecordRequest(BaseModel):
    meetingInfo: MeetingInfo
    transcript: str
    pptSummary: Optional[str] = ""
    papersSummary: Optional[str] = ""


@app.post("/generate-record")
def extract_pdf_text(file: UploadFile) -> str:
    try:
        reader = PdfReader(BytesIO(file.file.read()))
        texts = []
        for i, page in enumerate(reader.pages):
            txt = page.extract_text() or ""
            if txt.strip():
                texts.append(f"[PDF 第{i+1}页]\n{txt}")
        return "\n\n".join(texts)
    except Exception as e:
        return f"[PDF 解析失败: {str(e)}]"

def extract_ppt_text(file: UploadFile) -> str:
    try:
        prs = Presentation(BytesIO(file.file.read()))
        texts = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                texts.append(f"[PPT 第{i+1}页]\n" + "\n".join(slide_text))
        return "\n\n".join(texts)
    except Exception as e:
        return f"[PPT 解析失败: {str(e)}]"


def extract_docx_text(file: UploadFile) -> str:
    if Document is None:
        return "[Word 解析失败: 未安装 python-docx]"
    try:
        file_bytes = BytesIO(file.file.read())
        doc = Document(file_bytes)
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
        return "\n".join(texts)
    except Exception as e:
        return f"[Word 解析失败: {str(e)}]"


def extract_pdf_text_bytes(content: bytes) -> str:
    try:
        reader = PdfReader(BytesIO(content))
        texts = []
        for i, page in enumerate(reader.pages):
            txt = page.extract_text() or ""
            if txt.strip():
                texts.append(f"[PDF 第{i+1}页]\n{txt}")
        return "\n\n".join(texts)
    except Exception as e:
        return f"[PDF 解析失败: {str(e)}]"


def extract_pdf_pages_with_images(content: bytes) -> list:
    """按页提取 PDF：每页返回 { page, text, image }，image 为 base64 数据 URL。无 PyMuPDF 时仅返回 text。"""
    try:
        import fitz  # type: ignore  # PyMuPDF
    except ImportError:
        try:
            reader = PdfReader(BytesIO(content))
            return [
                {"page": i + 1, "text": (p.extract_text() or "").strip(), "image": None}
                for i, p in enumerate(reader.pages)
            ]
        except Exception:
            return []
    try:
        doc = fitz.open(stream=content, filetype="pdf")  # type: ignore
        pages = []
        for i in range(len(doc)):
            page = doc[i]
            text = (page.get_text() or "").strip()
            pix = page.get_pixmap(dpi=144)
            png_bytes = pix.tobytes("png")
            b64 = base64.b64encode(png_bytes).decode("ascii")
            pages.append({
                "page": i + 1,
                "text": text,
                "image": f"data:image/png;base64,{b64}",
            })
        doc.close()
        return pages
    except Exception:
        try:
            reader = PdfReader(BytesIO(content))
            return [
                {"page": i + 1, "text": (p.extract_text() or "").strip(), "image": None}
                for i, p in enumerate(reader.pages)
            ]
        except Exception:
            return []


def extract_docx_text_bytes(content: bytes) -> str:
    if Document is None:
        return "[Word 解析失败: 未安装 python-docx]"
    try:
        doc = Document(BytesIO(content))
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
        return "\n".join(texts)
    except Exception as e:
        return f"[Word 解析失败: {str(e)}]"


def extract_ppt_content_bytes(content: bytes):
    try:
        prs = Presentation(BytesIO(content))
        texts = []
        slides_meta = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            images = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shape.image
                    ext = (img.ext or "png").lower()
                    mime = "image/png"
                    if ext in ["jpg", "jpeg"]:
                        mime = "image/jpeg"
                    b64 = base64.b64encode(img.blob).decode("ascii")
                    images.append(f"data:{mime};base64,{b64}")
            if slide_text:
                texts.append(f"[PPT 第{i+1}页]\n" + "\n".join(slide_text))
            slides_meta.append(
                {
                    "page": i + 1,
                    "images": images,
                }
            )
        return "\n\n".join(texts), slides_meta
    except Exception as e:
        return f"[PPT 解析失败: {str(e)}]", []


def _find_soffice():
    """查找本机 soffice 可执行文件，优先 PATH，再常见路径（Windows / macOS / Linux）。"""
    exe = shutil.which("soffice") or shutil.which("libreoffice")
    if exe:
        return exe

    candidates = []

    # Windows 常见安装路径
    import sys
    if sys.platform == "win32":
        win_candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        # 兼容中文/英文系统的 Program Files 环境变量
        for env_var in ("PROGRAMFILES", "PROGRAMFILES(X86)", "PROGRAMW6432"):
            pf = os.environ.get(env_var)
            if pf:
                win_candidates.append(os.path.join(pf, "LibreOffice", "program", "soffice.exe"))
        candidates.extend(win_candidates)

    # macOS 常见路径
    candidates += [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/local/bin/soffice",
        "/opt/homebrew/bin/soffice",
    ]
    # macOS 模糊匹配（版本号不固定时）
    applications = Path("/Applications")
    if applications.exists():
        for app in applications.glob("LibreOffice*.app"):
            p = app / "Contents/MacOS/soffice"
            if p.exists():
                candidates.insert(0, str(p))

    for p in candidates:
        if os.path.isfile(p) and os.access(p, os.X_OK):
            return p
    return None


def _run_libreoffice_convert(ppt_path: Path, outdir: Path) -> bool:
    """使用 LibreOffice 将 PPT 转为 PDF，成功返回 True。"""
    ppt_abs = ppt_path.resolve()
    outdir_abs = outdir.resolve()
    base = ["--headless", "--convert-to", "pdf", str(ppt_abs), "--outdir", str(outdir_abs)]
    env = os.environ.copy()
    env.setdefault("HOME", str(outdir_abs))
    env["SAL_USE_VCLPLUGIN"] = "gen"
    candidates = []
    soffice_exe = _find_soffice()
    if soffice_exe:
        candidates.append([soffice_exe] + base)
    candidates.extend([
        ["libreoffice"] + base,
        ["soffice"] + base,
    ])
    for cmd in candidates:
        try:
            subprocess.run(
                cmd,
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=120,
                env=env,
                cwd=str(outdir_abs),
            )
            return True
        except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError):
            continue
    return False


def render_ppt_to_page_images(content: bytes):
    try:
        import fitz  # type: ignore
    except ImportError:
        return []

    # Windows 下 soffice/fitz 偶发短暂占用临时 PDF；Python 3.10+ 可用 ignore_cleanup_errors，3.9 用普通 TemporaryDirectory
    with TemporaryDirectory() as tmpdir:
        tmpdir_path = Path(tmpdir)
        ppt_path = tmpdir_path / "slides.pptx"
        pdf_path = tmpdir_path / "slides.pdf"
        ppt_path.write_bytes(content)

        if not _run_libreoffice_convert(ppt_path, tmpdir_path):
            return []

        if not pdf_path.exists():
            candidates = list(tmpdir_path.glob("*.pdf"))
            if not candidates:
                return []
            pdf_path = candidates[0]

        doc = None
        try:
            doc = fitz.open(pdf_path)  # type: ignore
        except Exception:
            return []

        images = []
        try:
            for page_index in range(len(doc)):
                page = doc[page_index]
                pix = page.get_pixmap(dpi=144)
                png_bytes = pix.tobytes("png")
                b64 = base64.b64encode(png_bytes).decode("ascii")
                images.append(f"data:image/png;base64,{b64}")
        finally:
            # Windows 下若不及时 close，临时目录清理会出现 WinError 32（文件占用）
            try:
                if doc is not None:
                    doc.close()
            except Exception:
                pass

        return images

@app.post("/process-meeting")
async def process_meeting(request: Request):
    form = await request.form()
    time = form.get("time") or ""
    if not time:
        from fastapi.responses import JSONResponse
        return JSONResponse(status_code=400, content={"detail": "缺少 time"})
    topic = form.get("topic") or ""
    roles = form.get("roles") or "[]"
    transcript = form.get("transcript") or ""
    meeting_type = (form.get("meeting_type") or "").strip()
    discipline = (form.get("discipline") or "").strip()
    files_list = form.getlist("files") or []
    files_list = [x for x in files_list if getattr(x, "filename", None)]
    def _parse_name_list(s: str) -> list:
        if not s or not s.strip():
            return []
        try:
            out = json.loads(s)
            if isinstance(out, list):
                return [str(x).strip() for x in out if str(x).strip()]
        except Exception:
            pass
        return [r.strip() for r in re.split(r"[,，、\s]+", s) if r.strip()]

    roles_list = _parse_name_list(roles)
    reporters_list = _parse_name_list(form.get("reporters") or "")

    ppt_texts = []
    ppt_slides = []
    pdf_texts = []
    papers_pages = []  # PDF 按页：{ page, text, image }，page 为全局 1-based
    transcript_extra_parts = []
    uploaded_file_names = []  # 原始文件名，供模型从文件名推断汇报人（如「张三汇报.pptx」）

    file_list = list(files_list) if files_list else []
    for f in file_list:
        original_filename = (getattr(f, "filename") or "").strip()
        uploaded_file_names.append(original_filename)
        name = original_filename.lower()
        try:
            if asyncio.iscoroutinefunction(getattr(f, "read", None)):
                content = await f.read()
            else:
                content = getattr(f, "file", None)
                content = content.read() if content else b""
        except Exception:
            content = b""
        if not isinstance(content, bytes):
            content = b""
        if name.endswith((".ppt", ".pptx")):
            text, slides_meta = extract_ppt_content_bytes(content)
            ppt_texts.append(text)
            page_images = render_ppt_to_page_images(content)
            if page_images:
                slides_meta = [
                    {"page": idx + 1, "images": [img]}
                    for idx, img in enumerate(page_images)
                ]
            ppt_slides.extend(slides_meta)
        elif name.endswith(".pdf"):
            pdf_texts.append(extract_pdf_text_bytes(content))
            pdf_pages = extract_pdf_pages_with_images(content)
            for p in pdf_pages:
                p["page"] = len(papers_pages) + 1
                papers_pages.append(p)
        elif name.endswith((".doc", ".docx")):
            docx_text = extract_docx_text_bytes(content)
            if any(kw in name for kw in ("转写", "语音", "转文字", "会议记录", "纪要", "实录", "记录")):
                transcript_extra_parts.append(docx_text)
            else:
                pdf_texts.append(docx_text)
        else:
            pass

    transcript_full = transcript.strip()
    if transcript_extra_parts:
        transcript_full = (transcript_full + "\n\n" + "\n\n".join(transcript_extra_parts)).strip()

    transcript_original = transcript_full

    ppt_combined = "\n\n".join(ppt_texts) if ppt_texts else ""
    pdf_combined = "\n\n".join(pdf_texts) if pdf_texts else ""

    # ── System Prompt ──────────────────────────────────────────────────────
    # 修改提示词请直接编辑下方。会议类型仅作参考，要点以本场会议实际讨论为准。
    # ──────────────────────────────────────────────────────────────────────

    MEETING_TYPE_GUIDES = {
        "人文社科": (
            "  1. \"核心论点\"：本次汇报的主要学术观点与论证脉络\n"
            "  2. \"论文结构与逻辑问题\"：写作结构、逻辑漏洞或老师指出的问题\n"
            "  3. \"文献综述要点\"：重要文献的讨论、引用建议与研究空白\n"
            "  4. \"修改意见与写作计划\"：明确的修改方向与下一步写作任务\n"
        ),
        "实验工科": (
            "  1. \"实验目的与方法\"：实验设计、研究问题与方法选择\n"
            "  2. \"数据与图表结果\"：关键实验数据、图表解读与主要结论\n"
            "  3. \"异常问题与原因分析\"：失败、异常结果及原因分析\n"
            "  4. \"下周实验安排\"：后续实验计划与参数调整\n"
        ),
        "理论数理": (
            "  1. \"研究思路与问题定义\"：核心研究问题与整体思路\n"
            "  2. \"公式/模型/算法要点\"：关键数学推导、模型结构或算法设计\n"
            "  3. \"当前难点\"：尚未解决的理论障碍或推导瓶颈\n"
            "  4. \"下一步推导方向\"：老师建议的后续理论探索方向\n"
        ),
        "通用组会": (
            "  1. \"本周完成情况\"：本周已完成的任务与进展\n"
            "  2. \"遇到的问题\"：本周遇到的困难、阻碍与未解决事项\n"
            "  3. \"老师意见与反馈\"：导师在会上给出的具体评价与建议\n"
            "  4. \"下周计划\"：下周的明确任务目标与时间节点\n"
        ),
    }

    type_guide_block = ""
    for t, guide in MEETING_TYPE_GUIDES.items():
        type_guide_block += f'\n\u25b7 meeting_type = "{t}" \u65f6\uff0c\u53ef\u4f9d\u6b21\u4f5c\u4e3a\u5173\u6ce8\u70b9\uff08\u4ec5\u53c2\u8003\uff09\uff1a\n{guide}'

    user_meeting_type = (meeting_type or "").strip()
    user_discipline = (discipline or "").strip()

    role_instruction = ""
    if user_meeting_type and user_meeting_type in MEETING_TYPE_GUIDES:
        role_instruction = (
            "\n【重要：以用户选择为准】用户已选择会议类型为「" + user_meeting_type + "」。\n"
            "- meeting_type 必须且只能使用该值，不要根据转写/PPT/文献内容自行推断或改写。\n"
            "- 请以该类型领域专家的身份撰写纪要，使用该领域常用术语与规范，使输出更具针对性和专业性。\n\n"
        )
    if user_discipline:
        role_instruction = (
            role_instruction
            + "\n【重要：以用户填写学科为准】用户指定了学科「" + user_discipline + "」。\n"
            "- 请以该学科领域专家身份撰写纪要，使用该学科常用术语、关注该领域典型问题与规范，使输出更具针对性和专业性。\n"
            "- 不要根据内容猜测学科或类型，全文须体现该学科视角。若用户未选择会议类型，meeting_type 可标为「通用组会」，但内容须紧扣该学科。\n\n"
        )
    if reporters_list:
        role_instruction = (
            role_instruction
            + "\n【重要：按用户指定汇报人展开，不得遗漏、不得删除】用户指定了本次汇报人名单：" + "、".join(reporters_list) + "。\n"
            "- by_reporter 必须严格按该名单顺序与人数输出：名单里有几人就输出几条，顺序与名单一致。\n"
            "- 二、核心讨论要点 与 三、导师反馈与行动项 均按每位汇报人分别展开；每条对应一名汇报人的 reporter、key_points、advisor_feedback、action_items。\n"
            "- 若某汇报人在素材中未识别到其汇报内容，该汇报人仍须保留一条，其 key_points、advisor_feedback、action_items 可为空数组，或写一条 key_points：title 为「未识别到该汇报人相关内容」、detail 为「（根据当前素材未匹配到该汇报人的发言或材料）」、evidence 为空数组。不得直接删掉该汇报人条目。\n"
            "- utterances 中 speaker 若能从转写中对应到上述姓名，请尽量使用用户提供的汇报人姓名，便于与 by_reporter 一致。\n\n"
        )

    system_prompt = (
        "【研行记 · 身份与定位】\n"
        "你是「研行记 (Research Action Note)」的核心引擎，面向科研组会场景，负责把语音转写、PPT、文献等多源素材加工成「高价值、可执行、可溯源」的组会纪要。\n"
        "用途：帮助师生在会后快速回顾讨论要点、导师反馈与行动项，并能在纪要中追溯到「谁在何时说了什么、对应哪一页材料」，从而减少信息遗漏、明确下一步任务。\n"
        "特点：多来源证据链（每条要点/反馈/行动项尽量挂接 transcript、PPT 页、文献片段）；按汇报人分块组织；导师说的话要细致保留、不合并不省略；会议类型仅作参考，具体表述以本场会议实际讨论内容为准。\n"
        "语气与风格：专业、克制、信息密度高。用学术化、正式的中文表述，不夸张、不抒情；重点放在「事实与建议」的准确传达，让读者读完即知本场组会讨论了什么、老师具体提了哪些意见、接下来要做什么。\n\n"

        "【输出形式】\n"
        "请严格按以下阶段工作，最终只输出一个 JSON 对象，不要输出中间思考过程。\n\n"
        + role_instruction +

        "【阶段0：会议类型】\n"
        "仅当用户未选择类型且未填写学科时，才根据转写、PPT、文献判断 meeting_type；否则必须以上方用户选择/填写的为准。\n"
        "类型说明（供参考）：\"人文社科\"多涉及论文写作、文献综述、理论框架；\"实验工科\"多涉及实验设计、数据与图表、工程排查；\"理论数理\"多涉及数学推导、公式建模、算法设计；\"通用组会\"为周期性进展汇报或难以归类的讨论。\n"
        "注意：会议类型只是标签与参考，不约束你如何组织要点。本场会议究竟讨论了什么，就写什么；要点条数、小标题可根据实际内容灵活增减，不必拘泥于类型下的固定几条。\n\n"

        "【阶段1：汇报人多源识别与语音角色梳理】\n"
        "汇报人可从以下多源综合识别，请结合使用：\n"
        "1) 用户填写的「汇报人」名单（若提供则优先，且 by_reporter 须严格按该名单展开、不删不漏）；\n"
        "2) 上传文件名：meeting_info 中的 uploaded_file_names 可能含汇报人姓名（如「张三汇报.pptx」「李四组会.pdf」），可从中推断谁汇报；\n"
        "3) PPT 正文：ppt_text 开头常为标题页或汇报人信息，可结合识别；\n"
        "4) 语音转写：将 transcript 拆分为 utterances，每条含 index、speaker、content，speaker 尽量区分「导师/老师」与汇报人；转写中出现的发言人姓名、称谓可作为汇报人依据。\n"
        "若出现多位学生/汇报人，不得合并为一人；speaker 命名需稳定，尽量与用户提供的汇报人姓名或文件名中的姓名一致。导师发言不得混入汇报人发言中。\n"
        "说明：系统会用 utterances 做纪要生成；溯源时仍展示用户提供的原始转写原文，便于对照当时原话。\n\n"

        "【阶段2：素材对齐与证据链构建（多来源）】\n"
        "每条 key_points、advisor_feedback、action_items 至少一条 evidence，尽量多来源（转写、PPT 页、文献片段）。\n"
        "evidence 含 type(transcript|ppt|paper)、location、quote、context、note。location 写清具体出处（如「语音转写」「PPT 第3页」「文献 XXX 第2节」）；不捏造页码或原文。\n"
        "重要：对 type=transcript 的 evidence，quote 必须为转写原文中的连续片段（可模糊对应），以便在溯源界面高亮显示并上下滚动查看上下文。\n\n"

        "【阶段3：按汇报人组织 + 要点结构】\n"
        "输出 by_reporter，每人含 reporter、key_points、advisor_feedback、action_items。\n"
        "若用户提供了汇报人名单：by_reporter 必须按该名单顺序、每人一条；未识别到内容的汇报人保留条目，key_points/advisor_feedback/action_items 留空或写「未识别到该汇报人相关内容」，不得删除。\n"
        "若用户未提供名单：从 transcript、文件名、PPT 等识别汇报人，by_reporter 覆盖所有识别到的汇报人，不得只输出部分。\n"
        "「本次组会要点」对应 key_points。下面按会议类型给了一些可参考的方面，仅供你组织时参考，不必严格照搬：可根据本场会议实际讨论了什么来定条数与标题。重点是把「本场会议真实讨论到的内容」写全、写细。\n"
        "每位有内容的汇报人至少输出 1 条 key_points；若某汇报人信息不足且用户未要求保留，可写“该汇报人本次有效信息较少”；若用户名单中有此人则必须保留其条目且不得删除。\n"
        + type_guide_block + "\n\n"

        "【导师/老师说的内容：务必细致、不遗漏】\n"
        "导师（或老师）在会上的发言是纪要的核心价值之一，必须细致呈现，不要漏点、不要合并成一句笼统话。\n"
        "- 每条独立的意见、建议、批评、肯定，尽量单独成条或单独成句写清，避免「老师提了几点意见」这种概括。\n"
        "- 老师提到的具体信息必须保留：文献名、作者、年份、书名、论文题目、方法名、工具名、数据来源、某章某节、某页某段等，在 advisor_feedback 的 content 或 action_items 的 description 中写清楚。例如：「导师建议阅读张三等(2020)《某某研究》中关于……的论述」「可参考 XX 方法/XX 数据集」「建议先做 A 再做 B」。\n"
        "- 老师指出的具体问题（如某处逻辑、某段表述、某组实验）要写清是「哪方面的问题」以及「老师的大致意见」，不要只写「老师认为需要修改」而丢失具体指向。\n"
        "- 老师给出的下一步任务、时间节点、交付物，要在 action_items 中写具体：做什么、做到什么程度、为什么（reason 可简要说明）。\n"
        "若转写或 PPT 中已有上述具体表述，请原样或适度整理后写入对应条目的 content/description，确保读者仅凭纪要就能还原老师的主要意见与要求。\n\n"

        "【行动项等级与表述】\n"
        "action_items 的 level：RED=必须完成、有明确期限或硬性要求；YELLOW=建议优化、应优先考虑；GREEN=探索性、可选尝试。description 写清「做什么」；reason 可简要写「为何被判定为该等级或依据哪句讨论」。\n\n"

        "【JSON 输出结构（必须含 by_reporter）】\n"
        "{\n"
        "  \"meeting_type\": \"人文社科 | 实验工科 | 理论数理 | 通用组会\",\n"
        "  \"basic_info\": { \"time\": \"...\", \"topic\": \"...\", \"roles\": [\"...\"] },\n"
        "  \"utterances\": [ { \"index\": 1, \"speaker\": \"导师\", \"content\": \"...\" } ],\n"
        "  \"by_reporter\": [\n"
        "    {\n"
        "      \"reporter\": \"汇报人姓名\",\n"
        "      \"key_points\": [ { \"title\": \"...\", \"detail\": \"...\", \"evidence\": [ { \"type\": \"...\", \"location\": \"...\", \"quote\": \"...\", \"context\": \"...\", \"note\": \"...\" } ] } ],\n"
        "      \"advisor_feedback\": [ { \"speaker\": \"导师\", \"content\": \"...\", \"evidence\": [ ... ] } ],\n"
        "      \"action_items\": [ { \"level\": \"RED|YELLOW|GREEN\", \"description\": \"...\", \"reason\": \"...\", \"evidence\": [ ... ] } ]\n"
        "    }\n"
        "  ],\n"
        "  \"summary\": \"150字内总结\"\n"
        "}\n\n"

        "【阶段4：表达与格式】\n"
        "- 几乎全部使用中文，必要时可以使用英文术语，措辞学术、正式、信息密度高。basic_info 从 meeting_info 拷贝或从素材中提炼。\n"
        "- by_reporter 与 summary 严格依据 transcript、ppt_text、papers_text 生成；不编造未出现的讨论；每条要点/反馈/行动项的 evidence 尽量覆盖多类来源。\n"
        "- 导师反馈与行动项中凡有具体文献、方法、问题指向、时间节点、交付要求的，一律写清，不省略、不泛化为「按老师意见修改」等笼统表述。\n"
    )

    user_content = {
        "meeting_info": {
            "time": time,
            "topic": topic,
            "roles": roles_list,
            "reporters": reporters_list,
            "uploaded_file_names": uploaded_file_names,
        },
        "user_meeting_type": user_meeting_type or None,
        "user_discipline": user_discipline or None,
        "transcript": transcript_full,
        "ppt_text": ppt_combined,
        "papers_text": pdf_combined,
    }

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": json.dumps(user_content, ensure_ascii=False)},
    ]

    response = client.chat.completions.create(
        model=MODEL_NAME,
        messages=messages,
        response_format={"type": "json_object"},
    )

    content = response.choices[0].message.content
    data = json.loads(content)

    data["model_used"] = MODEL_NAME
    if user_meeting_type and user_meeting_type in MEETING_TYPE_GUIDES:
        data["meeting_type"] = user_meeting_type

    data["raw_sources"] = {
        "transcript": transcript_full,
        "transcript_original": transcript_original,
        "ppt_text": ppt_combined,
        "papers_text": pdf_combined,
        "ppt_slides": ppt_slides,
        "papers_pages": papers_pages,
    }

    def ensure_evidence(obj, fallback_quote: str):
        if not isinstance(obj, dict):
            return
        ev = obj.get("evidence")
        if not ev or not isinstance(ev, list) or len(ev) == 0:
            snippet = (transcript_full or "")[:400].strip() or "(无转写)"
            obj["evidence"] = [{
                "type": "transcript",
                "location": "语音转写",
                "quote": (fallback_quote or obj.get("content") or obj.get("detail") or obj.get("description") or "")[:200] or "见上文",
                "context": snippet,
                "note": "根据转写整理",
            }]

    for block in data.get("by_reporter") or []:
        for kp in block.get("key_points") or []:
            ensure_evidence(kp, kp.get("detail") or kp.get("title"))
        for fb in block.get("advisor_feedback") or []:
            ensure_evidence(fb, fb.get("content"))
        for ai in block.get("action_items") or []:
            ensure_evidence(ai, ai.get("description"))

    for kp in data.get("key_points") or []:
        ensure_evidence(kp, kp.get("detail") or kp.get("title"))
    for fb in data.get("advisor_feedback") or []:
        ensure_evidence(fb, fb.get("content"))
    for ai in data.get("action_items") or []:
        ensure_evidence(ai, ai.get("description"))

    def _is_teacher_like(name: str) -> bool:
        if not name:
            return False
        n = str(name).strip().lower()
        teacher_keys = ["导师", "老师", "教授", "pi", "老板", "导师a", "导师b", "teacher"]
        return any(k in n for k in teacher_keys)

    def _norm_name(name: str) -> str:
        return "".join(str(name or "").strip().lower().split())

    def _block_matches_reporter(blk: dict, reporter_name: str) -> bool:
        if not blk or not isinstance(blk, dict):
            return False
        r = str(blk.get("reporter") or "").strip()
        if not r:
            return False
        return _norm_name(r) == _norm_name(reporter_name)

    # 若用户提供了汇报人名单：严格按名单顺序保证每人一条，不删不漏；未识别到内容的留空块
    by_blocks = data.get("by_reporter")
    if not isinstance(by_blocks, list):
        by_blocks = []

    if reporters_list:
        ordered_blocks = []
        used_indices = set()
        for user_reporter in reporters_list:
            user_reporter = str(user_reporter or "").strip()
            if not user_reporter:
                continue
            matched = None
            for i, blk in enumerate(by_blocks):
                if i in used_indices:
                    continue
                if _block_matches_reporter(blk, user_reporter):
                    matched = blk
                    used_indices.add(i)
                    break
            if matched is not None:
                ordered_blocks.append(matched)
            else:
                ordered_blocks.append({
                    "reporter": user_reporter,
                    "key_points": [{
                        "title": "未识别到该汇报人相关内容",
                        "detail": "（根据当前素材未匹配到该汇报人的发言或材料）",
                        "evidence": [],
                    }],
                    "advisor_feedback": [],
                    "action_items": [],
                })
        if len(ordered_blocks) < len([r for r in reporters_list if str(r).strip()]):
            seen = {_norm_name(str(b.get("reporter") or "")) for b in ordered_blocks}
            for user_reporter in reporters_list:
                user_reporter = str(user_reporter or "").strip()
                if not user_reporter or _norm_name(user_reporter) in seen:
                    continue
                ordered_blocks.append({
                    "reporter": user_reporter,
                    "key_points": [{
                        "title": "未识别到该汇报人相关内容",
                        "detail": "（根据当前素材未匹配到该汇报人的发言或材料）",
                        "evidence": [],
                    }],
                    "advisor_feedback": [],
                    "action_items": [],
                })
                seen.add(_norm_name(user_reporter))
        by_blocks = ordered_blocks
    else:
        existing_reporters = set()
        for blk in by_blocks:
            if not isinstance(blk, dict):
                continue
            r = str(blk.get("reporter") or "").strip()
            if r and (not _is_teacher_like(r)):
                existing_reporters.add(_norm_name(r))

        utterances = data.get("utterances") or []
        speaker_stats = {}
        for u in utterances:
            if not isinstance(u, dict):
                continue
            spk = str(u.get("speaker") or "").strip()
            content_u = str(u.get("content") or "").strip()
            if not spk or not content_u:
                continue
            if _is_teacher_like(spk):
                continue
            key = _norm_name(spk)
            if not key:
                continue
            if key not in speaker_stats:
                speaker_stats[key] = {"name": spk, "turns": 0, "chars": 0, "samples": []}
            speaker_stats[key]["turns"] += 1
            speaker_stats[key]["chars"] += len(content_u)
            if len(speaker_stats[key]["samples"]) < 2:
                speaker_stats[key]["samples"].append(content_u[:120])

        missing = []
        for key, st in speaker_stats.items():
            if (st["turns"] >= 2 or st["chars"] >= 60) and key not in existing_reporters:
                missing.append(st)

        for st in missing:
            sample_text = "；".join(st["samples"]).strip() or "该汇报人本次有效信息较少。"
            by_blocks.append({
                "reporter": st["name"],
                "key_points": [{
                    "title": "发言要点（自动补全）",
                    "detail": sample_text,
                    "evidence": [{
                        "type": "transcript",
                        "location": "语音转写",
                        "quote": sample_text[:120],
                        "context": (transcript_full or "")[:400] or sample_text,
                        "note": "根据转写自动补齐遗漏汇报人",
                    }]
                }],
                "advisor_feedback": [],
                "action_items": [],
            })

    if by_blocks:
        data["by_reporter"] = by_blocks

    chosen_type = (meeting_type or "").strip()
    chosen_discipline = (discipline or "").strip()
    if chosen_type in MEETING_TYPE_GUIDES:
        data["meeting_type"] = chosen_type
        data["user_selected_meeting_type"] = chosen_type
    if chosen_discipline:
        data["user_selected_discipline"] = chosen_discipline
        if chosen_type not in MEETING_TYPE_GUIDES:
            data["meeting_type"] = data.get("meeting_type") or "通用组会"

    return data
